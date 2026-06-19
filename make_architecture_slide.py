"""
Generate a one-page PowerPoint slide describing the application architecture.

Run:  python make_architecture_slide.py
Output: ARCHITECTURE.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# --- palette ---------------------------------------------------------------
INK = RGBColor(0x1F, 0x2A, 0x37)        # dark text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
USER = RGBColor(0x55, 0x5F, 0x6B)       # grey
ORCH = RGBColor(0x7C, 0x3A, 0xED)       # purple (orchestrator)
OLLAMA = RGBColor(0x16, 0xA3, 0x4A)     # green
WRAP = RGBColor(0x25, 0x63, 0xEB)       # blue
CLIENT = RGBColor(0x0E, 0x94, 0x88)     # teal
FGT = RGBColor(0xDC, 0x2A, 0x26)        # red (FortiGate brand-ish)
DOCS = RGBColor(0xCA, 0x8A, 0x04)       # gold (documentation)
CONF = RGBColor(0x6B, 0x72, 0x80)       # slate (config)
LINE = RGBColor(0x94, 0xA3, 0xB8)       # arrow grey

SW, SH = 13.333, 7.5
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


def box(x, y, w, h, fill, title, subtitle="", title_sz=12.5, sub_sz=8.5,
        fg=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = fill
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_top", "margin_bottom", "margin_left", "margin_right"):
        setattr(tf, m, Pt(2))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(title_sz)
    r.font.color.rgb = fg
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(sub_sz)
        r2.font.color.rgb = fg
    return sp


def arrow(x1, y1, x2, y2, label="", color=LINE, width=2.0, dash=False,
          double=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    le = cn.line._get_or_add_ln()
    le.append(le.makeelement(qn('a:tailEnd'), {'type': 'triangle'}))
    if double:
        le.append(le.makeelement(qn('a:headEnd'), {'type': 'triangle'}))
    if dash:
        le.append(le.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if label:
        lx, ly = (x1 + x2) / 2 - 0.85, (y1 + y2) / 2 - 0.16
        tb = slide.shapes.add_textbox(Inches(lx), Inches(ly), Inches(1.7), Inches(0.26))
        tf = tb.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = False
        tp = tf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        rr = tp.add_run()
        rr.text = label
        rr.font.size = Pt(8)
        rr.font.italic = True
        rr.font.color.rgb = INK
    return cn


def textbox(x, y, w, h, lines, size, bold=False, color=INK, bullet=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = ("•  " + txt) if bullet else txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


# --- title (compact) -------------------------------------------------------
tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.55), Inches(0.8))
tf = tb.text_frame
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
p = tf.paragraphs[0]
r = p.add_run()
r.text = "FortiGate Firewall AI Agent — Architecture"
r.font.size = Pt(24)
r.font.bold = True
r.font.color.rgb = ORCH
sp = tf.add_paragraph()
rs = sp.add_run()
rs.text = "Local, privacy-preserving LLM assistant — live FortiGate data + built-in docs, no cloud keys."
rs.font.size = Pt(10.5)
rs.font.color.rgb = INK

# --- nodes -----------------------------------------------------------------
ROW, H = 2.95, 0.95
midL = ROW + H / 2
# Ollama (top)
box(2.55, 1.30, 2.55, 0.8, OLLAMA, "Ollama  (local LLM)",
    "OLLAMA_BASE_URL · emits TOOL: calls")
# pipeline
box(0.45, ROW, 1.70, H, USER, "USER", "terminal Q & A")
box(2.55, ROW, 2.55, H, ORCH, "langchain_firewall_agent.py",
    "Orchestrator · prompt + tool loop")
box(5.45, ROW, 2.30, H, WRAP, "fortigate_api_wrapper.py",
    "Adapter · maps tools, formats")
box(8.05, ROW, 2.00, H, CLIENT, "fortigate.py", "REST client (httpx)")
box(10.35, ROW, 2.50, H, FGT, "FortiGate firewall", "FortiOS /api/v2")
# documentation + config (below)
box(2.55, 4.45, 2.55, 0.8, DOCS, "pdf_loader.py",
    "Built-in docs → prompt context")
box(5.45, 4.45, 2.30, 0.8, CONF, ".env  (config)",
    "host · token · log source · Ollama URL")

# --- arrows ----------------------------------------------------------------
arrow(2.15, midL, 2.55, midL, double=True)                         # user <-> orch
arrow(3.82, ROW, 3.82, 2.10, label="prompt / answer", double=True)  # orch <-> ollama
arrow(5.10, midL, 5.45, midL, label="TOOL: call")                  # orch -> wrapper
arrow(7.75, midL, 8.05, midL, label="method")                      # wrapper -> client
arrow(10.05, midL, 10.35, midL, label="HTTPS + token")             # client -> fgt
arrow(3.82, 4.45, 3.82, ROW + H, label="docs", color=DOCS, dash=True)   # docs -> orch
arrow(6.60, 4.45, 6.60, ROW + H, label="config", color=CONF, dash=True)  # .env -> wrapper

# --- flow strip ------------------------------------------------------------
textbox(0.45, 5.45, 12.45, 0.55,
        ["Traffic flow:  Question → prompt (system + docs) → LLM emits TOOL: call "
         "→ wrapper → client → HTTPS to FortiGate → JSON back → results "
         "accumulate, LLM re-prompted (≤3×) → answer + Recommendation + [PDF]."],
        size=10, bold=True)

# --- key points ------------------------------------------------------------
textbox(0.45, 6.05, 12.45, 1.30, [
    "Text-based tool protocol (model prints TOOL: name(args), parsed by regex) — model-agnostic.",
    "Context accumulates across iterations; real API errors surfaced (403/404), not swallowed.",
    "Log tools auto-fall-back disk → memory → forticloud; all settings via .env.",
], size=9.5, bullet=True)

prs.save("ARCHITECTURE.pptx")

# --- overflow check --------------------------------------------------------
EMU = 914400
overflow = []
for sh in slide.shapes:
    left, top = sh.left / EMU, sh.top / EMU
    right, bottom = left + sh.width / EMU, top + sh.height / EMU
    if right > SW + 0.01 or bottom > SH + 0.01 or left < -0.01 or top < -0.01:
        name = sh.text_frame.text[:30] if sh.has_text_frame else sh.shape_type
        overflow.append((name, round(right, 2), round(bottom, 2)))
if overflow:
    print("OVERFLOW detected (right/bottom in inches, slide is %sx%s):" % (SW, SH))
    for n, r, b in overflow:
        print("  -", n, "right=%s bottom=%s" % (r, b))
else:
    print("OK - all shapes within %.3f x %.1f in slide bounds" % (SW, SH))
print("Saved ARCHITECTURE.pptx")
